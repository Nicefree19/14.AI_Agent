#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Create UPDATED PPT: 센코어테크 제작현황 + P5 대시보드 데이터 통합."""
import sys, os, io

if sys.stdout.encoding and sys.stdout.encoding.lower().startswith("cp"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
TASK_DIR = os.path.join(_ROOT, "telegram_data", "tasks", "msg_52")
os.makedirs(TASK_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BLUE = RGBColor(0x1B, 0x36, 0x5F)
BLUE = RGBColor(0x2E, 0x5C, 0x9A)
LIGHT_BLUE = RGBColor(0x4A, 0x86, 0xC8)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
RED = RGBColor(0xCC, 0x33, 0x33)
GREEN = RGBColor(0x33, 0x99, 0x33)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_RED = RGBColor(0x99, 0x22, 0x22)
YELLOW_BG = RGBColor(0xFF, 0xF3, 0xCD)
RED_BG = RGBColor(0xF8, 0xD7, 0xDA)

def add_header_bar(slide, title_text):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

def add_footer(slide):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(7.1), prs.slide_width, Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "P5 복합동 | 센코어테크 제작현황 + 대시보드 통합 분석 | 2026-02-09~02-11 기준"
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.CENTER
    tf.margin_top = Inches(0.05)

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r_idx == len(data) - 1 and ('합계' in str(data[r_idx][0]) or '소계' in str(data[r_idx][0])):
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
            else:
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    return table

def add_kpi_box(slide, left, top, width, height, label, value, color=BLUE):
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

def add_priority_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """Table with color-coded priority rows."""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER if c_idx in [0, 3, 4] else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r_idx <= 2:  # Critical rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = RED_BG
            elif r_idx <= 4:  # High rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = YELLOW_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    return table

def add_bullet_list(slide, left, top, width, height, items, font_size=11):
    """Add a bullet list with (text, color, bold) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size + 1) if bold else Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(3)
    return tf


# ====================================================================
# SLIDE 1: Title (updated to reflect integrated analysis)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_text_box(slide, 1.5, 1.2, 10, 1.2,
    "P5 복합동 PSRC 제작 현황 분석", 40, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.5, 10, 0.6,
    "센코어테크 생산일보 + P5 대시보드 통합 보고서", 24, False, RGBColor(0xAA, 0xCC, 0xFF), PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 3.5, 10, 0.5,
    "생산일보 기준: 2026-02-09  |  대시보드 기준: W07 (2/5~2/10)", 16, False, RGBColor(0x88, 0xAA, 0xDD), PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10, 0.5,
    "원본: 원유엽 이사/센코어테크 → 류재호 전무 → 이동혁 소장", 14, False, RGBColor(0x88, 0xAA, 0xDD), PP_ALIGN.CENTER)

add_text_box(slide, 1.5, 5.3, 10, 0.5,
    "분석일시: 2026-02-14  |  총 9페이지 (제작현황 6p + 대시보드 3p)", 14, False, RGBColor(0x77, 0x99, 0xCC), PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 2: Executive Summary (KPI Dashboard) - UPDATED
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  Executive Summary - 주요 지표")
add_footer(slide)

# KPI boxes row 1
add_kpi_box(slide, 0.5, 1.2, 2.4, 1.3, "총 PSRC 계획", "3,886 PCS", BLUE)
add_kpi_box(slide, 3.2, 1.2, 2.4, 1.3, "1절 PSRC 진행률", "24.0%", LIGHT_BLUE)
add_kpi_box(slide, 5.9, 1.2, 2.4, 1.3, "앙카 생산율", "57.1%", GREEN)
add_kpi_box(slide, 8.6, 1.2, 2.4, 1.3, "전체 PSRC 진행률", "4.6%", ORANGE)
add_kpi_box(slide, 11.3, 1.2, 1.7, 1.3, "긴급 이슈", "5건", RED)

# Summary bullets
tf = add_text_box(slide, 0.5, 2.9, 5.8, 3.5, "", 12)
items = [
    ("긍정적 사항", GREEN, True),
    ("• 앙카 프레임: 57.1% 생산, 34% 출하 완료", DARK_GRAY, False),
    ("• 인주/우영/정인 앙카 100% 완료", DARK_GRAY, False),
    ("• 1절주 우영/인주/정인 30%+ 제작 진행", DARK_GRAY, False),
    ("• 원자재 93.5% 입고 완료", DARK_GRAY, False),
    ("• SHOP REV 일정 확정 (양재영 이사, 2/6)", DARK_GRAY, False),
    ("• 2절주 EP 스터드 타입 확정 (남궁승, 2/6)", DARK_GRAY, False),
]
for i, (text, color, bold) in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(12) if not bold else Pt(13)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(3)

tf2 = add_text_box(slide, 6.8, 2.9, 6.2, 3.5, "", 12)
issues = [
    ("주의 사항 (제작 + 대시보드)", RED, True),
    ("• C존(센코어) 1절주: 완전 미착수 (0/176)", DARK_GRAY, False),
    ("• 대성1/동명/오창: 1절 PSRC 미착수", DARK_GRAY, False),
    ("• FORM(압연롤) 약 450톤 미입고 (30%)", DARK_GRAY, False),
    ("• PTW 철근 과다 / FMCS 오류 / STUD 불량", DARK_GRAY, False),
    ("• PSRC Rev.3 검토 회신 긴급 (2/11 마감)", RED, False),
    ("• EP-105~108 상세 검토 필요 (SEN-668)", RED, False),
    ("• 담당자 미지정 669/671건 (100%)", ORANGE, False),
    ("• 앙카 HOLD 해제 조건 미확정 (SEN-097)", ORANGE, False),
]
for i, (text, color, bold) in enumerate(issues):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(12) if not bold else Pt(13)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(3)

# ====================================================================
# SLIDE 3: 1절주/2절주 PSRC 제작 현황
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  1절주 / 2절주 PSRC 제작 현황 (제작 현황 갑지)")
add_footer(slide)

table_data = [
    ["존 구분", "총수량", "HMB제작", "면조립", "대조립", "HMB+PSRC", "전기", "FORM", "엠베드", "도장", "출하"],
    ["A존(삼원)", "308", "187", "200", "200", "190", "151", "148", "135", "131", "110"],
    ["B존(동성)", "253", "77", "77", "66", "66", "57", "52", "48", "46", "33"],
    ["C존(센코어)", "176", "0", "0", "0", "0", "0", "0", "0", "0", "0"],
    ["소계", "737", "264", "277", "266", "256", "208", "200", "183", "177", "143"],
]
add_table(slide, 0.5, 1.2, 12.3, 2.5, 5, 11, table_data)

add_text_box(slide, 0.5, 3.9, 12, 0.4,
    "※ C존(센코어 설치구간) 1절주는 제작 실적이 전혀 없음 (0/176) - 주의 필요",
    12, True, RED)

add_text_box(slide, 0.5, 4.5, 12, 0.4,
    "2절주 현황 (총 737 PCS)", 15, True, DARK_BLUE)

table_data2 = [
    ["존 구분", "총수량", "HMB제작", "면조립", "대조립", "HMB+PSRC", "FORM", "엠베드", "도장", "출하"],
    ["A존(삼원)", "308", "72", "74", "74", "72", "0", "0", "0", "0"],
    ["B존(동성)", "253", "0", "0", "0", "0", "0", "0", "0", "0"],
    ["C존(센코어)", "176", "0", "0", "0", "0", "0", "0", "0", "0"],
    ["소계", "737", "72", "74", "74", "72", "0", "0", "0", "0"],
]
add_table(slide, 0.5, 5.0, 12.3, 2.0, 5, 10, table_data2)

# ====================================================================
# SLIDE 4: 집계표 - 업체별 공정현황
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  MEMBER LIST 집계표 - 업체별 공정현황")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 5, 0.4, "앙카(ANCHOR) 부문 - 총 737 PCS", 14, True, DARK_BLUE)

anchor_data = [
    ["업체", "수량", "FRAME조립", "제작율", "출하", "출하율"],
    ["인주", "55", "55", "100%", "55", "100%"],
    ["우영", "55", "55", "100%", "55", "100%"],
    ["정인", "44", "44", "100%", "44", "100%"],
    ["코엘", "358", "150", "42%", "53", "15%"],
    ["씨엔이엔지", "225", "117", "52%", "42", "19%"],
    ["합계", "737", "421", "57%", "249", "34%"],
]
add_table(slide, 0.5, 1.5, 5.8, 2.8, 7, 6, anchor_data)

add_text_box(slide, 6.8, 1.1, 6, 0.4, "1절 PSRC 부문 - 총 737 PCS", 14, True, DARK_BLUE)

psrc_data = [
    ["업체", "수량", "LC-FRAME", "COVER PL", "FORM", "엠베드", "제작율", "도장"],
    ["대성1", "28", "0", "0", "0", "0", "0%", "0"],
    ["동명", "95", "0", "0", "0", "0", "0%", "0"],
    ["오창", "50", "0", "0", "0", "0", "0%", "0"],
    ["우영", "174", "77", "73", "63", "54", "34.8%", "54"],
    ["인주", "206", "88", "83", "75", "67", "34.9%", "65"],
    ["정인", "184", "77", "72", "62", "62", "32.4%", "58"],
    ["합계", "737", "242", "228", "200", "183", "25.9%", "177"],
]
add_table(slide, 6.8, 1.5, 6.2, 3.0, 8, 8, psrc_data)

# ====================================================================
# SLIDE 5: 일일 생산 집계 & 원자재
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  일일 생산 집계 & 원자재 현황")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 5, 0.4, "1절 PSRC - 업체별 일일생산 (2/9 기준)", 14, True, DARK_BLUE)

daily_data = [
    ["업체", "계획", "전일누계", "금일생산", "생산누계", "생산율"],
    ["대성1", "28", "0", "0", "0", "0%"],
    ["동명", "95", "0", "0", "0", "0%"],
    ["오창", "50", "0", "0", "0", "0%"],
    ["우영", "174", "44", "10", "54", "31.0%"],
    ["인주", "206", "55", "10", "65", "31.6%"],
    ["정인", "184", "44", "14", "58", "31.5%"],
    ["소계", "737", "143", "34", "177", "24.0%"],
]
add_table(slide, 0.5, 1.5, 5.8, 3.0, 8, 6, daily_data)

add_text_box(slide, 0.5, 4.7, 5.5, 0.8,
    "앙카: 계획 737 | 전일 339 | 금일 82 | 누계 421 (57.1%)\n출하: 전일 187 | 금일 62 | 누계 249 (33.8%)",
    12, False, DARK_GRAY)

add_text_box(slide, 6.8, 1.1, 6, 0.4, "원자재 현황 (P5+P6 형강류, 2/6 기준)", 14, True, DARK_BLUE)

material_data = [
    ["품목", "청구(톤)", "입고(톤)", "가공(톤)", "입고율"],
    ["L-BAR", "1,140", "1,140", "575", "100%"],
    ["MAIN ANGLE", "3,534", "3,534", "1,724", "100%"],
    ["FORM(압연롤)", "1,488", "1,037", "580", "70%"],
    ["BH(BRK/후판)", "251", "251", "244", "100%"],
    ["철근", "461", "463", "50", "100%"],
    ["합계", "6,874", "6,424", "3,173", "93.5%"],
]
add_table(slide, 6.8, 1.5, 6.0, 2.5, 7, 5, material_data)

add_text_box(slide, 6.8, 4.2, 5.8, 0.4,
    "※ FORM(압연롤) 약 450톤 미입고 (30% 미납)", 12, True, RED)
add_text_box(slide, 6.8, 4.6, 5.8, 0.4,
    "※ 가공 진행률: 전체 약 49%", 12, False, ORANGE)

# ====================================================================
# SLIDE 6: 제작 이슈 정리 & 전무님 코멘트
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  제작 이슈 정리 & 전무님 코멘트")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 6, 0.4, "제작 이슈 (2026-02-05~06)", 16, True, RED)

issue_data = [
    ["No.", "일자", "이슈 내용", "조치 사항"],
    ["1", "02-05", "PTW 철근 과다로 제작 지연 발생", "현장 철근 매립 요청"],
    ["2", "02-06", "PTW 삭제 구간 페인트 처리 문의", "전부 페인트 요함 (위동표 팀장 확인)"],
    ["3", "02-06", "FMCS 오류로 실적 등록 불가", "강상민 프로 요청 중"],
    ["4", "02-06", "STUD 볼트 용접기 불량 발생", "신규 발주 필요"],
]
add_table(slide, 0.5, 1.5, 12.3, 2.2, 5, 4, issue_data,
          col_widths=[0.5, 0.8, 5.5, 5.0])

add_text_box(slide, 0.5, 4.0, 12, 0.4, "류재호 전무님 코멘트", 16, True, DARK_BLUE)

shape = slide.shapes.add_shape(5, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
shape.line.color.rgb = LIGHT_BLUE

tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)

comments = [
    ("원본 메일 (2026-02-11 15:25):", True, DARK_BLUE),
    ('"도표로된 첨부 파일 1절 2절 두장만 보면됩니다"', False, DARK_GRAY),
    ("", False, DARK_GRAY),
    ("답장 메일 (2026-02-11 15:54):", True, DARK_BLUE),
    ('"앞으로 양소장이랑 정보 공유 토록 원이사한테 이야기 해놓을께요"', False, DARK_GRAY),
    ("-> 이동혁 소장: \"자료 감사합니다\"", False, MEDIUM_GRAY),
]
for i, (text, bold, color) in enumerate(comments):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_after = Pt(3)


# ====================================================================
# SLIDE 7: 즉시 조치 항목 & Critical/High 이슈 (NEW - 대시보드)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  [대시보드] 즉시 조치 항목 & Critical/High 이슈")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 6, 0.4, "즉시 조치 항목 (W07 기준)", 14, True, RED)

action_data = [
    ["우선순위", "항목", "마감", "관련 SEN", "상태"],
    ["1 (긴급)", "PSRC 기둥 SHOP Rev.3 검토 회신", "2/11", "SEN-428", "긴급"],
    ["2 (긴급)", "EP-105~108 임베디드 플레이트 상세 검토", "2/11", "SEN-668", "진행중"],
    ["3 (높음)", "EP 변경에 따른 SHOP 영향 범위 확인", "금주", "SEN-428", "대기"],
    ["4 (높음)", "BCW/SCRUBBER 2B Zone SHOP 도면 추적", "금주", "SEN-495", "미접수"],
    ["5 (보통)", "오프닝 CAD 파일 정리 및 송부", "2/14", "SEN-544", "예정"],
]
add_priority_table(slide, 0.5, 1.5, 12.3, 2.5, 6, 5, action_data,
                   col_widths=[1.2, 5.5, 0.8, 1.2, 0.8])

add_text_box(slide, 0.5, 4.2, 6, 0.4, "Critical/High SEN 이슈 교차참조", 14, True, DARK_BLUE)

# Critical issues list
issues_list = [
    ("SEN-668: EP 변경 계산근거 요청 (Critical)", RED, True),
    ("  센구조 발생, EP-105~108 상세도 수신 완료, 구조 계산서 제출 필요", DARK_GRAY, False),
    ("", DARK_GRAY, False),
    ("SEN-428: EP 변경에 따른 SHOP 중단 위험 (High)", ORANGE, True),
    ("  EP 변경→PSRC SHOP Rev.3 발행, 2/11 검토 회신 필요", DARK_GRAY, False),
    ("", DARK_GRAY, False),
    ("SEN-097: 앙카볼트 HOLD 미해제 (High)", ORANGE, True),
    ("  HOLD 해제 조건 미확정, 기초→골조 공정 밀림 위험", DARK_GRAY, False),
    ("", DARK_GRAY, False),
    ("SEN-495: 효율화 SHOP 도면 미접수 (High)", ORANGE, True),
    ("  BCW/SCRUBBER 2B Zone 도면 미접수, SHOP 일정 누적 지연 가능", DARK_GRAY, False),
]
add_bullet_list(slide, 0.5, 4.6, 12.3, 2.4, issues_list, font_size=11)


# ====================================================================
# SLIDE 8: SHOP REV 일정표 (NEW - 대시보드)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  [대시보드] SHOP REV 일정표 (2/6 확정판)")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 5, 0.4, "1절 골조 SHOP", 14, True, DARK_BLUE)

shop1_data = [
    ["순서", "Zone", "도면 종류", "상태", "비고"],
    ["1", "1A1", "PSRC 기둥", "확정", "Rev.3 검토중"],
    ["2", "1A2", "PSRC 기둥", "확정", ""],
    ["3", "1B", "HMB 브라켓", "확정", ""],
    ["4", "1C", "PC 거더", "확정", ""],
]
add_table(slide, 0.5, 1.5, 5.5, 2.0, 5, 5, shop1_data,
          col_widths=[0.6, 0.7, 1.5, 0.8, 1.9])

add_text_box(slide, 6.8, 1.1, 6, 0.4, "2절 골조 SHOP", 14, True, DARK_BLUE)

shop2_data = [
    ["순서", "Zone", "도면 종류", "상태", "비고"],
    ["1", "2A", "PSRC 기둥", "확정", ""],
    ["2", "2B", "BCW/SCRUBBER", "미접수", "추적 필요"],
    ["3", "2C", "잔여 부재", "예정", ""],
]
add_table(slide, 6.8, 1.5, 6.0, 1.6, 4, 5, shop2_data,
          col_widths=[0.6, 0.7, 1.8, 0.8, 2.1])

add_text_box(slide, 0.5, 3.8, 6, 0.4, "EMBEDDED 순서", 14, True, DARK_BLUE)

embed_data = [
    ["순서", "대상", "비고"],
    ["1", "1A1 임배드", "일정 정정됨 (02/06)"],
    ["2", "1A2 임배드", "일정 정정됨 (02/06)"],
    ["3", "1B 임배드", ""],
    ["4", "2A 임배드", ""],
    ["5", "2B 임배드", ""],
    ["6", "2C 임배드", ""],
]
add_table(slide, 0.5, 4.2, 5.5, 2.5, 7, 3, embed_data,
          col_widths=[0.8, 2.0, 2.7])

# Right side: 확정 사항 + 미결 의사결정
add_text_box(slide, 6.8, 3.8, 6, 0.4, "금주 확정 사항", 14, True, GREEN)

confirmed_data = [
    ["날짜", "확정 내용", "확정자"],
    ["02/06", "2절주 EP 스터드 타입 최종 확정", "남궁승"],
    ["02/06", "1A1/1A2 임배드 일정 정정", "양재영 이사"],
]
add_table(slide, 6.8, 4.2, 6.0, 1.0, 3, 3, confirmed_data,
          col_widths=[0.8, 3.5, 1.7])

add_text_box(slide, 6.8, 5.4, 6, 0.4, "미결 의사결정", 14, True, ORANGE)

pending_data = [
    ["항목", "영향 범위", "긴급도"],
    ["TC 좌대 변경 승인", "기초 -> 골조", "보통"],
    ["SS-Splice 적용 여부", "PSRC 전체", "보통"],
    ["앙카 HOLD 해제 조건", "기초 공사", "긴급"],
]
add_table(slide, 6.8, 5.8, 6.0, 1.2, 4, 3, pending_data,
          col_widths=[2.5, 1.8, 1.7])


# ====================================================================
# SLIDE 9: 데이터 품질 경고 & 다음 주 업무 (NEW - 대시보드)
# ====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "  [대시보드] 데이터 품질 경고 & 다음 주 업무 계획")
add_footer(slide)

add_text_box(slide, 0.5, 1.1, 5, 0.4, "데이터 품질 경고", 14, True, RED)

quality_data = [
    ["수준", "항목", "현황", "영향"],
    ["긴급", "담당자 미지정", "669/671건 (100%)", "책임 추적 불가"],
    ["긴급", "마감일 없는 H/C", "522건", "우선순위 정렬 불가"],
    ["주의", "결정 기록 미비", "1/671건 (0.1%)", "이력 추적 불가"],
    ["주의", "리뷰 큐 적체", "14건 미처리", "이슈 흐름 정체"],
]
add_priority_table(slide, 0.5, 1.5, 5.8, 2.0, 5, 4, quality_data,
                   col_widths=[0.7, 1.5, 1.6, 2.0])

# Recommendation box
shape = slide.shapes.add_shape(5, Inches(0.5), Inches(3.7), Inches(5.8), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "권장 조치"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = GREEN
p2 = tf.add_paragraph()
p2.text = "Critical/High 이슈 5건에 우선적으로 담당자/마감일 지정"
p2.font.size = Pt(11)
p2.font.color.rgb = DARK_GRAY
p3 = tf.add_paragraph()
p3.text = "-> Google Sheets 양방향 동기화로 이해관계자 실시간 공유"
p3.font.size = Pt(11)
p3.font.color.rgb = DARK_GRAY

# Email analysis summary
add_text_box(slide, 0.5, 5.0, 5.8, 0.4, "금주 이메일 분석 (25건)", 14, True, DARK_BLUE)

email_data = [
    ["카테고리", "건수", "주요 내용"],
    ["EP 관련", "8건", "11열 상세도, 스터드 타입 확정"],
    ["PSRC/SHOP", "7건", "Rev.3 검토, 일정 확정, 출도 현황"],
    ["기초/구조", "2건", "철근 배근, 앙카 HOLD"],
    ["기타 업무", "2건", "오프닝 CAD, 주간회의"],
    ["비P5", "6건", "일반 공지, 타 프로젝트"],
]
add_table(slide, 0.5, 5.4, 5.8, 1.5, 6, 3, email_data,
          col_widths=[1.2, 0.8, 3.8])

# Right side: Next week plan
add_text_box(slide, 6.8, 1.1, 6, 0.4, "다음 주 예상 업무 (2/11~2/14)", 14, True, DARK_BLUE)

next_week_items = [
    ("긴급 (2/11 마감)", RED, True),
    ("  1. PSRC Rev.3 검토 회신 (삼우)", DARK_GRAY, False),
    ("  2. EP-105~108 상세 검토 (스터드/앵커 간섭)", DARK_GRAY, False),
    ("  3. EP 계산근거 작성/송부 (SEN-668 대응)", DARK_GRAY, False),
    ("", DARK_GRAY, False),
    ("금주 중 (2/14까지)", ORANGE, True),
    ("  4. BCW/SCRUBBER 도면 추적 (2B Zone)", DARK_GRAY, False),
    ("  5. 오프닝 CAD 파일 정리 (2/14 마감)", DARK_GRAY, False),
    ("  6. 앙카 HOLD 해제 조건 협의 (센구조+ENA)", DARK_GRAY, False),
    ("", DARK_GRAY, False),
    ("일반 업무", DARK_BLUE, True),
    ("  7. SHOP 출도 현황 업데이트", DARK_GRAY, False),
    ("  8. 주간회의 안건 준비", DARK_GRAY, False),
    ("  9. SS-Splice 검토 결과 정리", DARK_GRAY, False),
    ("  10. OCR 교정 검토 (저확신도 도면번호)", DARK_GRAY, False),
]
add_bullet_list(slide, 6.8, 1.5, 6.0, 5.5, next_week_items, font_size=11)


# ====================================================================
# SAVE
# ====================================================================
output_path = os.path.join(TASK_DIR, "P5_센코어테크_제작현황_통합분석.pptx")
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Size: {os.path.getsize(output_path):,} bytes")
print(f"Slides: {len(prs.slides)}")
print("SUCCESS")
